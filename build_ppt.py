"""Build the ClassGraph presentation.

Deck is structured 1:1 against the five evaluation criteria. Every number in
the results/problems slides is a real measured value from this repository's
own test runs -- see CHALLENGES_AND_SOLUTIONS.md for provenance.

Run:  python build_ppt.py
Out:  ClassGraph.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "ppt_assets"
OUT = ROOT / "ClassGraph.pptx"

# --------------------------------------------------------------------------- #
# Design system
# --------------------------------------------------------------------------- #

SW, SH = 13.333, 7.5          # slide size, inches (16:9)
ML = 0.68                      # left margin
MR = 0.68                      # right margin
CW = SW - ML - MR              # content width

INK = RGBColor(0x10, 0x27, 0x3F)        # deep navy - headings
BODY = RGBColor(0x3C, 0x50, 0x66)       # slate - body text
MUTE = RGBColor(0x7B, 0x8C, 0x9E)       # muted grey
TEAL = RGBColor(0x0E, 0x7C, 0x86)       # primary accent
TEAL_D = RGBColor(0x08, 0x59, 0x61)
AMBER = RGBColor(0xB4, 0x7A, 0x14)      # attention / problem
GREEN = RGBColor(0x1F, 0x6F, 0x50)      # solved / done
RED = RGBColor(0xA9, 0x33, 0x2A)        # blocker
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF3, 0xF7, 0xFA)      # card fill
PANEL2 = RGBColor(0xE9, 0xF1, 0xF4)     # tinted card
BORDER = RGBColor(0xDA, 0xE3, 0xEB)
RULE = RGBColor(0xE6, 0xEC, 0xF2)

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"
MONO = "Consolas"


def P(text, size=12, bold=False, color=BODY, font=FONT, italic=False,
      align=None, space_before=None, space_after=None, line=None):
    """One paragraph with a single run."""
    return {
        "runs": [{"t": text, "size": size, "bold": bold, "color": color,
                  "font": font, "italic": italic}],
        "align": align, "space_before": space_before,
        "space_after": space_after, "line": line,
    }


def PR(runs, align=None, space_before=None, space_after=None, line=None):
    """One paragraph with multiple runs. `runs` = list of dicts."""
    return {"runs": runs, "align": align, "space_before": space_before,
            "space_after": space_after, "line": line}


def R(t, size=12, bold=False, color=BODY, font=FONT, italic=False):
    return {"t": t, "size": size, "bold": bold, "color": color,
            "font": font, "italic": italic}


def add_text(slide, x, y, w, h, blocks, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, blk in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if blk.get("align") is not None:
            p.alignment = blk["align"]
        if blk.get("space_before") is not None:
            p.space_before = Pt(blk["space_before"])
        if blk.get("space_after") is not None:
            p.space_after = Pt(blk["space_after"])
        if blk.get("line") is not None:
            p.line_spacing = blk["line"]
        for seg in blk["runs"]:
            r = p.add_run()
            r.text = seg["t"]
            f = r.font
            f.name = seg.get("font", FONT)
            f.size = Pt(seg.get("size", 12))
            f.bold = seg.get("bold", False)
            f.italic = seg.get("italic", False)
            f.color.rgb = seg.get("color", BODY)
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.045):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = adj
        except (IndexError, KeyError):
            pass
    s.text_frame.word_wrap = True
    return s


def bar(slide, x, y, w, h, fill):
    """Flat rectangle, no line -- for rules and accent bars."""
    return rect(slide, x, y, w, h, fill=fill, line=None, shape=MSO_SHAPE.RECTANGLE)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _wrapped_height(text, size, width_in, line=1.26):
    """Estimate rendered height so headers never silently sit on the content
    below them. Segoe UI averages ~0.50em glyph advance for mixed-case text."""
    chars_per_line = max(1.0, (width_in * 72.0) / (0.50 * size))
    import math as _m
    lines = max(1, _m.ceil(len(text) / chars_per_line))
    return lines * size * line / 72.0


def chrome(prs, criterion, title, lead=None, page=None):
    """Standard slide header + footer. Returns (slide, content_top).

    Title and lead boxes are sized from their own estimated wrapped height, so
    a two-line title pushes the content down instead of overlapping it.
    """
    s = new_slide(prs)
    y = 0.46
    if criterion:
        bar(s, ML, y + 0.045, 0.11, 0.16, TEAL)
        add_text(s, ML + 0.24, y, 8.0, 0.22,
                 [P(criterion.upper(), 10.5, True, TEAL, FONT_SB)])
        y += 0.30
    th = _wrapped_height(title, 27, CW) + 0.06
    add_text(s, ML, y, CW, th, [P(title, 27, True, INK, FONT_SB)])
    y += th + 0.11
    if lead:
        lh = _wrapped_height(lead, 12.5, CW - 0.2, line=1.18) + 0.04
        add_text(s, ML, y, CW - 0.2, lh, [P(lead, 12.5, False, MUTE, line=1.15)])
        y += lh + 0.10
    if page is not None:
        footer(s, page)
    return s, y + 0.10


def footer(s, page):
    bar(s, ML, SH - 0.62, CW, 0.012, RULE)
    add_text(s, ML, SH - 0.50, 6.0, 0.24,
             [P("ClassGraph", 9, False, MUTE)])
    add_text(s, SW - MR - 1.2, SH - 0.50, 1.2, 0.24,
             [P(str(page), 9, True, MUTE, align=PP_ALIGN.RIGHT)])


def card(slide, x, y, w, h, heading=None, lines=None, accent=TEAL,
         fill=PANEL, heading_size=13, body_size=10.5, accent_bar=True):
    """Panel with an optional left accent bar, heading and body lines."""
    rect(slide, x, y, w, h, fill=fill)
    if accent_bar:
        bar(slide, x, y + 0.10, 0.055, h - 0.20, accent)
    tx = x + (0.26 if accent_bar else 0.20)
    tw = w - (tx - x) - 0.20
    blocks = []
    if heading:
        blocks.append(P(heading, heading_size, True, INK, FONT_SB,
                        space_after=5, line=1.05))
    for ln in (lines or []):
        if isinstance(ln, str):
            blocks.append(P(ln, body_size, False, BODY, space_after=3, line=1.18))
        else:
            blocks.append(ln)
    if blocks:
        add_text(slide, tx, y + 0.16, tw, h - 0.30, blocks)


def stat(slide, x, y, w, h, value, label, color=TEAL):
    rect(slide, x, y, w, h, fill=PANEL)
    add_text(slide, x + 0.14, y + 0.16, w - 0.28, 0.60,
             [P(value, 30, True, color, FONT_SB, align=PP_ALIGN.CENTER)])
    add_text(slide, x + 0.14, y + 0.82, w - 0.28, h - 0.95,
             [P(label, 9.5, False, BODY, align=PP_ALIGN.CENTER, line=1.12)])


def table(slide, x, y, w, col_w, data, row_h=0.34, head_h=0.36,
          head_fill=INK, head_color=WHITE, size=9.5, head_size=9.5,
          zebra=True, col_colors=None, col_bold=None):
    """Simple, fully-styled table. data[0] is the header row."""
    rows, cols = len(data), len(data[0])
    total_h = head_h + row_h * (rows - 1)
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                Inches(w), Inches(total_h))
    tbl = gt.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cwid in enumerate(col_w):
        tbl.columns[i].width = Inches(cwid)
    for r in range(rows):
        tbl.rows[r].height = Inches(head_h if r == 0 else row_h)
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = head_fill
            else:
                cell.fill.fore_color.rgb = (
                    WHITE if (not zebra or r % 2 == 1) else PANEL
                )
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            f = run.font
            f.name = FONT_SB if r == 0 else FONT
            f.size = Pt(head_size if r == 0 else size)
            f.bold = bool(r == 0 or (col_bold and c in col_bold))
            if r == 0:
                f.color.rgb = head_color
            elif col_colors and c in col_colors:
                f.color.rgb = col_colors[c]
            else:
                f.color.rgb = BODY
    return tbl


def pill(slide, x, y, w, h, text, fill, color=WHITE, size=9):
    rect(slide, x, y, w, h, fill=fill, line=None, adj=0.5)
    add_text(slide, x, y + 0.035, w, h - 0.05,
             [P(text, size, True, color, FONT_SB, align=PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #

def s01_title(prs):
    s = new_slide(prs)
    bar(s, 0, 0, SW, 0.10, TEAL)
    bar(s, 0, SH - 0.055, SW, 0.055, INK)

    add_text(s, ML, 1.62, CW, 0.4,
             [P("CV COURSE PROJECT", 11.5, True, TEAL, FONT_SB)])
    add_text(s, ML, 2.05, CW, 1.05, [P("ClassGraph", 58, True, INK, FONT_SB)])
    add_text(s, ML, 3.16, 9.6, 0.62,
             [P("Temporal Scene-Graph Group-Activity Analytics for Classrooms",
                20, False, TEAL_D, FONT_SB, line=1.1)])
    bar(s, ML, 3.95, 1.5, 0.028, TEAL)
    add_text(s, ML, 4.20, 9.3, 1.0,
             [P("Identity-persistent, behaviour-level engagement analytics — turning raw "
                "classroom video into per-student attention trajectories and group-level "
                "insight, without reducing students to a single score.",
                13, False, BODY, line=1.32)])

    labels = [("Perception + Identity", "Stages 1–2 built & tested"),
              ("Scene Graph + Temporal", "Stages 3–4 designed"),
              ("Group Activity", "Stage 5 designed")]
    bw = (CW - 0.44) / 3
    for i, (h, sub) in enumerate(labels):
        x = ML + i * (bw + 0.22)
        card(s, x, 5.55, bw, 0.86, heading=h, lines=[sub],
             heading_size=11.5, body_size=9.5, accent=TEAL if i == 0 else MUTE)
    return s


def s02_roadmap(prs, page):
    s, y = chrome(prs, None, "Review Roadmap",
                  "This deck is organised 1:1 against the five evaluation criteria.",
                  page)
    items = [
        ("1", "Problem Understanding", "Clarity of problem, objectives, motivation", 3),
        ("2", "Research Paper Selection", "Relevance, novelty, correctness of base paper", 6),
        ("3", "Literature Review Depth", "Related work identification & comparison", 7),
        ("4", "Dataset Identification", "Appropriateness, data format, challenges", 10),
        ("5", "Project Feasibility & Plan", "Workflow, timeline, results, risks", 11),
    ]
    lw = 7.35
    for i, (num, title, desc, slide_no) in enumerate(items):
        yy = y + 0.06 + i * 0.80
        rect(s, ML, yy, lw, 0.68, fill=WHITE, line=BORDER)
        rect(s, ML + 0.14, yy + 0.13, 0.42, 0.42, fill=INK, line=None, adj=0.22)
        add_text(s, ML + 0.14, yy + 0.19, 0.42, 0.30,
                 [P(num, 13, True, WHITE, FONT_SB, align=PP_ALIGN.CENTER)])
        add_text(s, ML + 0.70, yy + 0.11, 4.0, 0.26,
                 [P(title, 12.5, True, INK, FONT_SB)])
        add_text(s, ML + 0.70, yy + 0.36, 5.2, 0.24,
                 [P(desc, 9.5, False, MUTE)])
        add_text(s, ML + lw - 1.15, yy + 0.22, 1.0, 0.26,
                 [P(f"slide {slide_no}", 9, True, TEAL, align=PP_ALIGN.RIGHT)])

    px = ML + lw + 0.30
    pw = CW - lw - 0.30
    rect(s, px, y + 0.06, pw, 3.34, fill=INK, line=None)
    add_text(s, px + 0.30, y + 0.32, pw - 0.60, 0.3,
             [P("THE THESIS", 10, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])
    add_text(s, px + 0.30, y + 0.72, pw - 0.60, 2.4,
             [P("Represent a classroom as a dynamic scene graph — students are nodes, "
                "attention behaviours are edges — keep identity stable within the session, "
                "and reason over how that graph evolves to recognise individual and group "
                "engagement.", 13.5, False, WHITE, line=1.34)])
    card(s, px, y + 3.56, pw, 0.86,
         heading="Status at this review",
         lines=["Stages 1–2 implemented, 109 automated tests passing, "
                "validated on real classroom footage."],
         heading_size=11.5, body_size=9.5, accent=GREEN, fill=PANEL2)
    return s


def s03_problem(prs, page):
    s, y = chrome(prs, "Criterion 1 — Problem Understanding",
                  "Attendance is solved. Engagement is not.",
                  "The classroom already has cameras. What it does not have is any way to "
                  "answer the question a teacher actually cares about.", page)

    card(s, ML, y, 6.15, 1.62, heading="The problem",
         lines=["A teacher cannot continuously observe 30–40 students while also teaching.",
                "Manual observation is subjective, sampled at a few moments, and does not "
                "scale to every student across a whole lecture."],
         accent=RED)
    card(s, ML + 6.45, y, 6.15, 1.62, heading="Why existing tools fall short",
         lines=["Attendance and face-recognition systems answer “who is present”, "
                "not “who is engaged”.",
                "Emotion classifiers give a per-frame label with no memory — no system "
                "reasons about behaviour over time or across the group."],
         accent=RED)

    yy = y + 1.86
    add_text(s, ML, yy, CW, 0.3,
             [P("What we mean by “engagement” — defined before it is measured",
                13.5, True, INK, FONT_SB)])
    add_text(s, ML, yy + 0.34, CW, 1.00,
             [PR([R("We measure ", 11), R("observable behaviour", 11, True, TEAL),
                  R(", not inner emotional state. ", 11),
                  R("Where a student is looking, whether their eyes are open, whether a "
                    "phone is present, how their body is oriented, and how all of that "
                    "changes over time. ", 11),
                  R("We never infer feelings — that is both scientifically weaker and, in "
                    "education settings, legally restricted.", 11, False, AMBER)],
                 line=1.22)])

    yy += 1.46
    objectives = [
        ("Detect", "every student in frame, including back rows"),
        ("Track", "stable identity for the whole session"),
        ("Estimate", "attention state over a time window"),
        ("Flag", "device distraction & eye closure"),
        ("Relate", "peer orientation & group activity"),
        ("Report", "trajectories + class-level trends"),
    ]
    bw = (CW - 5 * 0.16) / 6
    for i, (verb, desc) in enumerate(objectives):
        x = ML + i * (bw + 0.16)
        rect(s, x, yy, bw, 1.12, fill=PANEL2, line=BORDER)
        add_text(s, x + 0.14, yy + 0.14, bw - 0.28, 0.26,
                 [P(verb, 12, True, TEAL_D, FONT_SB)])
        add_text(s, x + 0.14, yy + 0.44, bw - 0.28, 0.62,
                 [P(desc, 9, False, BODY, line=1.16)])
    return s


def s04_motivation(prs, page):
    s, y = chrome(prs, "Criterion 1 — Motivation",
                  "Why this matters, and why now",
                  None, page)

    # Left: two calm, real stats. Right: the sharper argument -- deployed
    # tools already exist, and their design is the mistake we are avoiding.
    lw = 4.55
    stat(s, ML, y, lw, 1.62, "30–40",
         "students one teacher must watch at once in a typical class", color=TEAL)
    stat(s, ML, y + 1.86, lw, 1.62, "~50%",
         "of detected students have no usable face from a real ceiling camera "
         "— measured on our own footage", color=TEAL)

    rx = ML + lw + 0.28
    rw = CW - lw - 0.28
    rect(s, rx, y, rw, 3.66, fill=INK, line=None)
    add_text(s, rx + 0.30, y + 0.20, rw - 0.60, 0.26,
             [P("THIS PROBLEM IS ALREADY BEING SOLVED BADLY", 10, True,
                RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])

    add_text(s, rx + 0.30, y + 0.50, 1.1, 0.50,
             [P("83", 32, True, WHITE, FONT_SB)])
    add_text(s, rx + 1.30, y + 0.54, rw - 1.60, 0.70,
             [P("Hong Kong schools now using “4 Little Trees”, which reads a "
                "child's face and scores their emotion in real time.",
                11.5, False, WHITE, line=1.22)])
    bar(s, rx + 0.30, y + 1.34, rw - 0.60, 0.012, RGBColor(0x2E, 0x47, 0x5F))

    add_text(s, rx + 0.30, y + 1.52, rw - 0.60, 0.86,
             [PR([R("A Hangzhou school", 12, True, WHITE, FONT_SB),
                  R(" scans every student's face every 30 seconds and posts a live "
                    "attention score on the classroom wall.", 12, False, WHITE)],
                 line=1.26)])
    bar(s, rx + 0.30, y + 2.46, rw - 0.60, 0.012, RGBColor(0x2E, 0x47, 0x5F))

    add_text(s, rx + 0.30, y + 2.62, rw - 0.60, 0.90,
             [PR([R("Both were built on facial emotion recognition and both drew real "
                    "privacy backlash", 11.5, True, AMBER),
                  R(" — the exact design mistake this project deliberately avoids.",
                    11.5, False, WHITE)], line=1.28)])

    yy = y + 3.86
    card(s, ML, yy, (CW - 0.28) / 2, 0.96, heading="Why now — feasibility changed",
         lines=["Detection, face mesh, head pose and body pose are all strong pretrained "
                "models today. The open question is no longer “can we see it” but “what "
                "does it mean, and over what timescale”."],
         accent=TEAL, heading_size=11, body_size=9)
    card(s, ML + (CW - 0.28) / 2 + 0.28, yy, (CW - 0.28) / 2, 0.96,
         heading="Where it applies",
         lines=["Schools, colleges and coaching centres · hybrid learning · corporate "
                "training · teaching-quality feedback · educational research on attention."],
         accent=TEAL, heading_size=11, body_size=9)

    yy += 1.14
    rect(s, ML, yy, CW, 0.82, fill=PANEL2, line=TEAL, line_w=1.1)
    add_text(s, ML + 0.28, yy + 0.14, CW - 0.56, 0.60,
             [PR([R("Our framing:  ", 11.5, True, TEAL_D, FONT_SB),
                  R("behaviour, not emotion — and a trend a teacher can act on, not a "
                    "score per student. “Attention across the room dropped sharply 22 "
                    "minutes in” is actionable and fair. A live per-child score on the "
                    "classroom wall is neither.", 11.5, False, INK)], line=1.22)])
    return s


def s05_objectives(prs, page):
    s, y = chrome(prs, "Criterion 1 — Objectives",
                  "Objectives, each with a measurable success criterion",
                  "Stated so progress can be checked rather than asserted. "
                  "Status reflects this review.", page)
    data = [
        ["#", "Objective", "Measurable success criterion", "Status"],
        ["O1", "Detect every student in frame",
         "Recall on real classroom frames; back rows not lost to downscaling", "Done"],
        ["O2", "Maintain stable identity within a session",
         "One track_id per student across occlusion; identity never crosses sessions", "Done"],
        ["O3", "Per-student attention state over time",
         "Windowed distribution over 6 behaviour categories, not per-frame verdicts", "Done"],
        ["O4", "Flag device distraction and eye closure",
         "Phone-overlap + gaze-down evidence; EAR-based eye state", "Done"],
        ["O5", "Model peer orientation and group activity",
         "Geometric pair detection now; scene-graph + GAR next", "In progress"],
        ["O6", "Report trajectories and class-level trends",
         "Class summary by default, individual as explicit drill-down", "Partial"],
    ]
    tbl = table(s, ML, y, CW, [0.52, 3.30, 6.55, 1.60], data,
                row_h=0.52, head_h=0.40, size=10, head_size=10,
                col_bold={1})
    # colour the status column by value
    status_colors = {"Done": GREEN, "In progress": AMBER, "Partial": AMBER}
    for r in range(1, len(data)):
        cell = tbl.cell(r, 3)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = status_colors.get(data[r][3], BODY)
        run.font.bold = True

    add_text(s, ML, y + 3.60, CW, 0.5,
             [PR([R("Note on scope:  ", 10.5, True, INK, FONT_SB),
                  R("O1–O4 are implemented and covered by automated tests. O5–O6 are "
                    "partially built — the honest position at this review is a working "
                    "perception + attention core with the relational layer designed and "
                    "started, not finished.", 10.5)], line=1.2)])
    return s


def s06_papers(prs, page):
    s, y = chrome(prs, "Criterion 2 — Research Paper Selection",
                  "Three base papers: robust temporal facial-signal modelling + group activity",
                  "All three are peer-reviewed, address a real sub-problem we face, and are "
                  "used for their method, not their task — see the note below on why.", page)
    papers = [
        (TEAL, "TEMPORAL ROBUSTNESS · 1 OF 2",
         "Robust Dynamic Facial Expression Recognition",
         "Liu, Wang & Shen — IEEE Trans. on Biometrics, Behavior, and Identity Science, 2025",
         "A dual-stream network that disentangles short-term facial movements from "
         "longer-term state, and explicitly separates genuinely hard samples from "
         "noisy/mislabelled ones using prediction agreement across resampled clips.",
         "What we take: NOT emotion classification. We take the principle of "
         "separating a momentary signal from a sustained state — our 15-second "
         "rolling window does exactly this for gaze and posture."),
        (TEAL, "TEMPORAL ROBUSTNESS · 2 OF 2",
         "Dynamic Objectives Learning for Facial Expression Recognition",
         "Wen, Chang, Li & Jiang — IEEE Trans. on Multimedia, vol. 22, pp. 2914–2925, 2020",
         "Splits training into stages with different objectives, and proposes a loss "
         "that explicitly widens the gap between easily-confused expression categories "
         "rather than forcing every sample into a single confident class.",
         "What we take: the same lesson without training a classifier — our "
         "“head-down, no device” category stays deliberately separate and uncertain "
         "instead of being forced into “engaged” or “distracted”."),
        (TEAL, "GROUP-LEVEL REASONING",
         "Learning Actor Relation Graphs for Group Activity Recognition",
         "Wu, Wang, Wang, Guo & Wu — CVPR 2019",
         "Builds an Actor Relation Graph capturing both appearance and position "
         "relations between actors, learned end-to-end with a GCN. Benchmarked on "
         "Volleyball and Collective Activity.",
         "What we take: the actor-relation graph + GCN readout to turn per-student "
         "states into one class-level activity label — used directly, not adapted."),
    ]
    cwid = (CW - 2 * 0.24) / 3
    card_h = 4.00
    for i, (accent, tag, title, cite, method, take) in enumerate(papers):
        x = ML + i * (cwid + 0.24)
        rect(s, x, y, cwid, card_h, fill=WHITE, line=BORDER)
        bar(s, x, y, cwid, 0.055, accent)
        add_text(s, x + 0.22, y + 0.20, cwid - 0.44, 0.24,
                 [P(tag, 8.5, True, TEAL, FONT_SB)])
        add_text(s, x + 0.22, y + 0.48, cwid - 0.44, 0.90,
                 [P(title, 12.5, True, INK, FONT_SB, line=1.10)])
        add_text(s, x + 0.22, y + 1.42, cwid - 0.44, 0.46,
                 [P(cite, 8.5, True, TEAL_D, line=1.18)])
        add_text(s, x + 0.22, y + 1.94, cwid - 0.44, 1.08,
                 [P(method, 9.5, False, BODY, line=1.20)])
        bar(s, x + 0.22, y + 3.06, cwid - 0.44, 0.012, RULE)
        add_text(s, x + 0.22, y + 3.16, cwid - 0.44, 0.80,
                 [P(take, 9, True, TEAL_D, line=1.18)])

    rect(s, ML, y + card_h + 0.14, CW, 0.80, fill=PANEL2, line=BORDER)
    add_text(s, ML + 0.26, y + card_h + 0.28, CW - 0.52, 0.56,
             [PR([R("Why facial-expression papers for a system that doesn't classify emotion:  ",
                    11, True, INK, FONT_SB),
                  R("we deliberately excluded emotion inference (see slide 3), but the "
                    "*methodology* both papers use for handling ambiguous, noisy facial "
                    "signal over time is exactly our problem too — we cite the technique, "
                    "not the task.", 11)], line=1.22)])
    return s


def s07_lit_technical(prs, page):
    s, y = chrome(prs, "Criterion 3 — Literature Review Depth",
                  "Related work, method, result, and the gap each leaves open",
                  "Not just our 3 base papers — the actual state of the art in each pillar, "
                  "so the gap we claim is checked against what already exists, not assumed.",
                  page)
    data = [
        ["Area", "Key work (venue, year)", "Method in one line", "Benchmark", "Limitation for our setting"],
        ["Cloth-changing\nRe-ID",
         "CAL (CVPR 2022);\nAIM, DeepChange (2023)",
         "Adversarial / causal loss suppresses clothing cues to keep RGB-only identity",
         "PRCC, LTCC,\nDeepChange",
         "Built for surveillance re-identification; never applied to session-scoped classroom identity"],
        ["Dynamic/robust\nfacial signal",
         "Liu et al. (2025);\nWen et al. (2020)",
         "Disentangle short-term movement from sustained state; separate hard samples "
         "from noisy ones",
         "DFEW, FERV39K",
         "Built for emotion classification, not behaviour geometry — we take the "
         "method (temporal + ambiguity handling), not the task"],
        ["Group activity\nrecognition",
         "ARG (CVPR 2019);\nGroupFormer (ICCV 2021)",
         "Actor relation graph / joint spatio-temporal attention over a GCN",
         "Volleyball,\nCollective Activity",
         "Sports/crowd domains; one group label per clip, not sustained engagement"],
        ["Classroom CV",
         "SCB-Dataset (2023–25);\nDAiSEE, EmotiW",
         "Per-frame behaviour detection; affect classification from webcam video",
         "SCB, DAiSEE",
         "No identity across time, no relational or group reasoning, single-subject framing"],
    ]
    table(s, ML, y, CW, [1.35, 2.55, 3.35, 1.55, 3.17], data,
          row_h=0.82, head_h=0.36, size=9.5, head_size=9.5, col_bold={0})

    add_text(s, ML, y + 3.66, CW, 0.60,
             [PR([R("Reading across the rows:  ", 10.5, True, INK, FONT_SB),
                  R("identity, relations and group reasoning have each been solved "
                    "separately, and classroom CV has solved per-frame detection. "
                    "No row does two of these at once — that gap is what we build on next.",
                    10.5)], line=1.2)])
    return s


def s08_lit_behavioural(prs, page):
    s, y = chrome(prs, "Criterion 3 — Literature Review Depth",
                  "The second half: what should the labels even be?",
                  "A CV-only review would stop at the previous slide. But the hardest question "
                  "here is not architectural — it is what counts as “engaged”, and over what "
                  "timescale. We reviewed the education and psychology literature to answer it.",
                  page)
    boxes = [
        ("1. Most systems invent their own labels", RED,
         "A field-wide review found most engagement-detection systems make up their own "
         "categories instead of using a validated one — so their results can't even be "
         "compared to each other. We took that as a warning about our own design."),
        ("2. So we used a real, validated one instead", GREEN,
         "BOSS is a published classroom-observation method already used by school "
         "psychologists: on-task vs. three kinds of off-task. Our six categories are "
         "built to match it, and we say plainly where they can't."),
        ("3. We picked our time window from real studies", TEAL,
         "One study showed a short mental break actually protects attention rather than "
         "harming it. Two separate lines of research — gaze tracking and classroom "
         "observation — both independently point to measuring over about 10–15 seconds, "
         "not judging a single frame."),
        ("4. We know where a camera has to stop guessing", AMBER,
         "Two students turned toward each other could be discussing the lesson or just "
         "chatting — telling those apart needs hearing them, not just seeing them. "
         "So we report “these two are facing each other” and refuse to guess why."),
    ]
    colw = (CW - 0.30) / 2
    bh = 1.56
    for i, (heading, accent, body) in enumerate(boxes):
        col, row = i % 2, i // 2
        x = ML + col * (colw + 0.30)
        yy = y + row * (bh + 0.22)
        card(s, x, yy, colw, bh, heading=heading, lines=[body],
             accent=accent, heading_size=12.5, body_size=10.5)

    yy2 = y + 2 * (bh + 0.22) + 0.06
    rect(s, ML, yy2, CW, 0.72, fill=INK, line=None)
    add_text(s, ML + 0.28, yy2 + 0.12, CW - 0.56, 0.50,
             [PR([R("Bottom line:  ", 11, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB),
                  R("the 15-second window already running in our code came from this "
                    "research, not a guess — and two unrelated fields agree on it.",
                    11, False, WHITE)], line=1.22)])
    return s


def s09_gap(prs, page):
    s, y = chrome(prs, "Criterion 3 — Research Gap",
                  "The white space, stated precisely",
                  None, page)
    cols = [
        ("Identity persists", "Re-ID gives a stable student across a session",
         "but only ever demonstrated for surveillance"),
        ("Relations are modelled", "VidSGG gives evolving subject–predicate–object graphs",
         "but with generic predicates, never attention"),
        ("The group is read out", "GAR gives one activity label for a set of actors",
         "but on sports clips, not sustained engagement"),
    ]
    cwid = (CW - 2 * 0.26) / 3
    for i, (h, has, lacks) in enumerate(cols):
        x = ML + i * (cwid + 0.26)
        rect(s, x, y, cwid, 1.86, fill=WHITE, line=BORDER)
        add_text(s, x + 0.22, y + 0.22, cwid - 0.44, 0.3,
                 [P(h, 13, True, INK, FONT_SB)])
        add_text(s, x + 0.22, y + 0.62, cwid - 0.44, 0.56,
                 [PR([R("✓ ", 11, True, GREEN), R(has, 10, False, BODY)], line=1.18)])
        add_text(s, x + 0.22, y + 1.24, cwid - 0.44, 0.52,
                 [PR([R("✗ ", 11, True, RED), R(lacks, 10, False, MUTE)], line=1.18)])

    yy = y + 2.14
    rect(s, ML, yy, CW, 1.30, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.30, yy + 0.20, CW - 0.60, 0.30,
             [P("THE GAP WE ADDRESS", 10, True, TEAL_D, FONT_SB)])
    add_text(s, ML + 0.30, yy + 0.56, CW - 0.60, 0.62,
             [P("No prior work combines identity persistence, relational scene graphs and "
                "temporal group reasoning for classrooms — and none of the classroom systems "
                "ground their behaviour labels in a validated observation instrument.",
                14, False, INK, line=1.26)])

    yy += 1.58
    card(s, ML, yy, CW, 1.06,
         heading="Our claim is deliberately narrow",
         lines=["We are not claiming a new Re-ID loss, a new scene-graph architecture, or a "
                "new GAR backbone. The contribution is the integration and the domain: an "
                "identity-persistent, behaviour-grounded, temporally-reasoned classroom "
                "analytics system — plus the honest reporting of where vision alone cannot "
                "decide, which the reviewed classroom literature consistently glosses over."],
         accent=TEAL, heading_size=12, body_size=10.5)
    return s


def s10_dataset_primary(prs, page):
    s, y = chrome(prs, "Criterion 4 — Dataset Identification",
                  "Primary dataset: OUC-CGE (Classroom Group Engagement)",
                  "Reselected after the previous candidate turned out to be a moving target: "
                  "an arXiv-only preprint revised seven times with no fixed, citable version.",
                  page)

    card(s, ML, y, 4.05, 3.10, heading="Why this one",
         lines=["Peer-reviewed and published in Nature Scientific Data (2025) — a fixed, "
                "citable dataset paper, not a preprint that keeps changing under us.",
                "Filmed with three real classroom cameras, including one mounted above the "
                "teacher's head — the same elevated angle behind our hardest problem.",
                "Labels are group-level engagement (High / Medium / Low), which is exactly "
                "the unit our Group Activity stage (Pillar 3, ARG) needs to be evaluated on.",
                "Openly archived with a permanent DOI on OSF, MIT-licensed code — one stable "
                "link, not a repository that has been rewritten under the same name."],
         accent=GREEN, heading_size=13, body_size=9.5)

    x2 = ML + 4.05 + 0.26
    rect(s, x2, y, 3.55, 3.10, fill=PANEL, line=BORDER)
    add_text(s, x2 + 0.24, y + 0.20, 3.1, 0.3,
             [P("Scale & format", 13, True, INK, FONT_SB)])
    rows = [("Video segments", "~7,705 clips"),
            ("Total footage", "12h 50m"),
            ("Participants", "17 students"),
            ("Camera angles", "3 (incl. overhead)"),
            ("Label unit", "group, per clip"),
            ("Archive", "OSF · DOI, MIT code")]
    for i, (k, v) in enumerate(rows):
        yy = y + 0.62 + i * 0.42
        add_text(s, x2 + 0.24, yy, 1.85, 0.28, [P(k, 9.5, False, MUTE)])
        add_text(s, x2 + 2.05, yy, 1.28, 0.28,
                 [P(v, 9.5, True, INK, MONO, align=PP_ALIGN.RIGHT)])

    x3 = x2 + 3.55 + 0.26
    w3 = SW - MR - x3
    rect(s, x3, y, w3, 3.10, fill=INK, line=None)
    add_text(s, x3 + 0.26, y + 0.20, w3 - 0.52, 0.3,
             [P("BEING HONEST ABOUT ITS LIMITS", 10, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])
    add_text(s, x3 + 0.26, y + 0.56, w3 - 0.52, 0.62,
             [P("Small (17 students, 16 female/1 male) and group-level only.",
                12.5, True, WHITE, FONT_SB, line=1.20)])
    add_text(s, x3 + 0.26, y + 1.24, w3 - 0.52, 1.66,
             [P("It cannot tell us whether ONE student is bowed over reading vs. bowed over "
                "a phone — that per-student behaviour granularity is what the supporting "
                "dataset below is for. We are using OUC-CGE for what it is actually strong "
                "at: real classroom camera angles and citable group-engagement ground truth.",
                10, False, RGBColor(0xD6, 0xE4, 0xEC), line=1.26)])

    yy = y + 3.34
    card(s, ML, yy, CW, 1.06,
         heading="Supporting dataset — SCB-Dataset, for per-student behaviour labels",
         lines=["OUC-CGE has no equivalent to this: SCB's bow head / read / write / using "
                "the phone / turn head / talk / discuss classes are supervision for the exact "
                "ambiguity our own geometry cannot resolve (see Problem 6). Used as a "
                "secondary reference, not the primary dataset, precisely because of the "
                "version-stability issue above — cite specific frozen commits, not \"latest\"."],
         accent=AMBER, heading_size=11.5, body_size=9.5, fill=PANEL)
    return s


def s12_architecture(prs, page):
    s, y = chrome(prs, "Criterion 5 — Workflow & Plan",
                  "Five stages, one team of three, and where we are today",
                  "Each stage consumes the previous one's output through a fixed JSON "
                  "schema, so three people build in parallel against stable interfaces. "
                  "Stage number doubles as the project month.", page)
    stages = [
        ("1", "Perception", "YOLOv11 · Face Mesh\nBody pose · head pose",
         "who & where", "A", GREEN, "BUILT"),
        ("2", "Identity", "ByteTrack now\ncloth-invariant Re-ID next",
         "stable IDs", "A", GREEN, "BUILT"),
        ("3", "Scene Graph", "nodes = students\nedges = behaviours",
         "what relates to what", "B", MUTE, "NEXT"),
        ("4", "Temporal", "sequence decoder\nover the graph",
         "how it evolves", "B", MUTE, "NEXT"),
        ("5", "Group Activity", "ARG-style graph\nreadout + dashboard",
         "the class as a whole", "C", MUTE, "NEXT"),
    ]
    gap = 0.28
    bw = (CW - 4 * gap) / 5
    for i, (num, title, tech, out, owner, accent, badge) in enumerate(stages):
        x = ML + i * (bw + gap)
        done = badge == "BUILT"
        rect(s, x, y, bw, 2.86, fill=WHITE if done else PANEL, line=accent if done else BORDER,
             line_w=1.5 if done else 0.75)
        bar(s, x, y, bw, 0.05, accent)
        add_text(s, x + 0.16, y + 0.16, 1.0, 0.26,
                 [P(f"MONTH {num}", 8.5, True, MUTE, FONT_SB)])
        pw = 0.62
        pill(s, x + bw - pw - 0.16, y + 0.16, pw, 0.24, badge,
             GREEN if done else RGBColor(0xC3, 0xCD, 0xD6), WHITE, size=7)
        add_text(s, x + 0.16, y + 0.50, bw - 0.32, 0.32,
                 [P(title, 12.5, True, INK, FONT_SB)])
        add_text(s, x + 0.16, y + 0.88, bw - 0.32, 0.82,
                 [P(tech, 8.5, False, BODY, MONO, line=1.22)])
        bar(s, x + 0.16, y + 1.80, bw - 0.32, 0.012, RULE)
        add_text(s, x + 0.16, y + 1.92, bw - 0.32, 0.40,
                 [P(out, 9, True, TEAL_D, line=1.14)])
        rect(s, x + 0.16, y + 2.36, 0.36, 0.36, fill=INK, line=None, adj=0.22)
        add_text(s, x + 0.16, y + 2.42, 0.36, 0.26,
                 [P(owner, 11, True, WHITE, FONT_SB, align=PP_ALIGN.CENTER)])
        add_text(s, x + 0.58, y + 2.42, bw - 0.74, 0.28,
                 [P("owner", 8, False, MUTE)])

    yy = y + 3.10
    card(s, ML, yy, (CW - 0.26) / 2, 1.10,
         heading="The frozen contract + privacy by design",
         lines=["One JSON object per frame, validated against a fixed schema on every run. "
                "Identity is computed from motion/geometry only, scoped to one session — "
                "no face-recognition database — and two tests enforce it cannot leak "
                "between sessions."],
         accent=TEAL, heading_size=11.5, body_size=9.5)
    card(s, ML + (CW - 0.26) / 2 + 0.26, yy, (CW - 0.26) / 2, 1.10,
         heading="The risk that is already retired",
         lines=["The biggest open question — “can this actually work on real classroom "
                "footage from a bad camera angle?” — has been answered empirically, with "
                "seven real failure modes found and fixed (next two slides)."],
         accent=GREEN, heading_size=11.5, body_size=9.5)
    return s


def s14_results(prs, page):
    s, y = chrome(prs, "Criterion 5 — Results",
                  "Measured improvements, before and after",
                  "Every figure below is from a real run on real classroom images in this "
                  "repository — not an estimate.", page)
    data = [
        ["Metric", "Before", "After", "Change"],
        ["Persons detected — busiest classroom image", "17", "19", "+2"],
        ["Faces detected on real footage — 12 images", "0", "95", "0 → usable"],
        ["Persons on a 640×480 webcam frame", "0", "1", "demo unblocked"],
        ["Gaze label “down” reachable at all", "No (bug)", "Yes", "fixed"],
        ["Automated tests passing", "10 (9 skipped)", "423 (0 skipped)", "+413"],
        ["Throughput, 4K video on RTX 4050", "11.0 FPS", "7.8 FPS", "cost of posture"],
    ]
    tw = 8.20
    tbl = table(s, ML, y, tw, [4.40, 1.25, 1.25, 1.30], data,
                row_h=0.42, head_h=0.38, size=9.5, head_size=9.5, col_bold={0})
    for r in range(1, len(data)):
        for c in (1, 2, 3):
            run = tbl.cell(r, c).text_frame.paragraphs[0].runs[0]
            run.font.name = MONO
            run.font.size = Pt(9)
            if c == 3:
                # last row is a cost, not a win -- colour it honestly
                run.font.color.rgb = AMBER if r == len(data) - 1 else GREEN
                run.font.bold = True
            elif c == 1:
                run.font.color.rgb = MUTE

    yy = y + 0.38 + 5 * 0.42 + 0.26
    card(s, ML, yy, tw, 1.06,
         heading="Why the throughput row is on this slide",
         lines=["Adding the body-pose fallback took us from 11.0 to 7.8 FPS. We are reporting "
                "the cost next to the benefit rather than quoting only the coverage gain — "
                "the trade is worth it for offline analysis, and it is the number a reviewer "
                "should be able to challenge us on."],
         accent=AMBER, heading_size=11.5, body_size=9.5, fill=PANEL)

    img_w = 2.05   # assets are square (dataset is 640x640); two stacked must fit
    img_x = SW - MR - img_w
    a = ASSETS / "annot_baseline_960_c40.jpg"
    b = ASSETS / "annot_candidate_1536_c30.jpg"
    if a.exists() and b.exists():
        # Measure the rendered asset rather than assuming its shape: the images
        # are regenerated from the dataset, whose aspect ratio is not fixed.
        from PIL import Image
        with Image.open(a) as _im:
            ih = img_w * _im.height / _im.width
        add_text(s, img_x, y, img_w, 0.24,
                 [P("BEFORE — 960 px, conf 0.40 → 17 found", 9, True, RED, FONT_SB)])
        s.shapes.add_picture(str(a), Inches(img_x), Inches(y + 0.28),
                             width=Inches(img_w))
        y2 = y + 0.28 + ih + 0.20
        add_text(s, img_x, y2, img_w, 0.24,
                 [P("AFTER — 1536 px, conf 0.30 → 19 found", 9, True, GREEN, FONT_SB)])
        s.shapes.add_picture(str(b), Inches(img_x), Inches(y2 + 0.28),
                             width=Inches(img_w))
        add_text(s, img_x, y2 + 0.28 + ih + 0.10, img_w, 0.54,
                 [P("Both panels are the same frame, rendered by tools/make_ppt_assets.py "
                    "from the current code — so the count under each picture is the count "
                    "the shipped pipeline produces, not a figure typed in by hand.",
                    8, False, MUTE, line=1.16)])
    return s


def s15_problems_perception(prs, page):
    s, y = chrome(prs, "Criterion 5 — Problems Faced (1 / 3)",
                  "Perception problems, root causes, and what we did",
                  "These were found by running the pipeline on real footage — none of them "
                  "were visible in synthetic tests.", page)

    probs = [
        ("PROBLEM 1", "Camera angle — a hard ceiling, not a bug",
         "Real classroom cameras sit high in a rear corner. A student bowed over a desk "
         "shows the camera the crown of their head. We measured a ~45% ceiling on face "
         "availability: 130 of 236 detected students had no face that ANY face model "
         "could find.",
         "Accepted it as a constraint instead of chasing it. Added a face-independent "
         "signal (body pose) so those students are still measurable — coverage went from "
         "265/321 to 321/321.", RED),
        ("PROBLEM 2", "Face detection returned zero faces on real footage",
         "Face Mesh was run once over the whole frame. MediaPipe downscales its input, so "
         "a face that is small relative to a 4K frame is destroyed before detection runs. "
         "Result: 0 faces on every real image, while synthetic tests passed.",
         "Run Face Mesh per detected person crop instead. 0/139 → 95/139 faces. We also "
         "measured that padding the crop makes it worse, and shipped zero padding against "
         "our own initial intuition.", RED),
    ]
    LB = 8.40                      # left block: the two detailed problems
    img_w = CW - LB - 0.30         # right column: the fallback, working
    cwid = (LB - 0.28) / 2
    for i, (tag, title, cause, fix, accent) in enumerate(probs):
        x = ML + i * (cwid + 0.28)
        rect(s, x, y, cwid, 3.30, fill=WHITE, line=BORDER)
        bar(s, x, y, cwid, 0.05, accent)
        add_text(s, x + 0.24, y + 0.20, cwid - 0.48, 0.24,
                 [P(tag, 9, True, accent, FONT_SB)])
        add_text(s, x + 0.24, y + 0.48, cwid - 0.48, 0.52,
                 [P(title, 13, True, INK, FONT_SB, line=1.08)])
        add_text(s, x + 0.24, y + 1.04, cwid - 0.48, 1.12,
                 [P(cause, 10, False, BODY, line=1.20)])
        bar(s, x + 0.24, y + 2.22, cwid - 0.48, 0.012, RULE)
        add_text(s, x + 0.24, y + 2.34, cwid - 0.48, 0.92,
                 [PR([R("Solution:  ", 10, True, GREEN, FONT_SB),
                      R(fix, 10, False, TEAL_D)], line=1.20)])

    # The posture fallback is the answer to Problem 1 -- show it working.
    pose = ASSETS / "t1_pose.jpg"
    if pose.exists():
        px = ML + LB + 0.30
        add_text(s, px, y, img_w, 0.24,
                 [P("THE FALLBACK, WORKING", 8.5, True, GREEN, FONT_SB)])
        from PIL import Image
        with Image.open(pose) as _im:
            ih = img_w * _im.height / _im.width
        s.shapes.add_picture(str(pose), Inches(px), Inches(y + 0.28),
                             width=Inches(img_w))
        add_text(s, px, y + 0.28 + ih + 0.10, img_w, 0.90,
                 [P("Body pose recovered for students with no detectable face. The single "
                    "upright skeleton is the standing teacher — correctly distinguished "
                    "from every seated student, which is a real check that the signal "
                    "means something.", 8.5, False, MUTE, line=1.18)])

    yy = y + 3.54
    small = [("PROBLEM 3", "Back-row students missed entirely",
              "Inference resolution 960 px shrank a 60 px student to ~30 px.",
              "imgsz 1280 + conf 0.30 → 139 to 236 persons."),
             ("PROBLEM 4", "Low-resolution input collapses face detection",
              "Measured 7% face rate below 480 px vs 44% at 720–1080 px.",
              "Set a minimum capture resolution as a deployment requirement.")]
    cwid2 = (CW - 0.28) / 2
    for i, (tag, title, cause, fix) in enumerate(small):
        x = ML + i * (cwid2 + 0.28)
        rect(s, x, yy, cwid2, 1.42, fill=PANEL, line=BORDER)
        add_text(s, x + 0.22, yy + 0.16, cwid2 - 0.44, 0.22,
                 [P(tag, 8.5, True, AMBER, FONT_SB)])
        add_text(s, x + 0.22, yy + 0.40, cwid2 - 0.44, 0.30,
                 [P(title, 10.5, True, INK, FONT_SB, line=1.10)])
        add_text(s, x + 0.22, yy + 0.74, cwid2 - 0.44, 0.34,
                 [P(cause, 9, False, BODY, line=1.16)])
        add_text(s, x + 0.22, yy + 1.10, cwid2 - 0.44, 0.26,
                 [PR([R("→ ", 9, True, GREEN), R(fix, 9, False, TEAL_D)], line=1.16)])

    return s


def s16_problems_interpretation(prs, page):
    s, y = chrome(prs, "Criterion 5 — Problems Faced (2 / 3)",
                  "Interpretation problems — and where we say “we cannot tell”",
                  None, page)

    # Problem 5 - gaze inversion
    rect(s, ML, y, 6.05, 2.30, fill=WHITE, line=BORDER)
    bar(s, ML, y, 6.05, 0.05, RED)
    add_text(s, ML + 0.24, y + 0.22, 5.6, 0.24,
             [P("PROBLEM 5", 9, True, RED, FONT_SB)])
    add_text(s, ML + 0.24, y + 0.50, 5.6, 0.30,
             [P("Gaze direction was silently inverted", 13, True, INK, FONT_SB)])
    add_text(s, ML + 0.24, y + 0.90, 5.6, 0.80,
             [P("The head-pose model reports pitch as up-positive; our contract assumed "
                "down-positive. Students bowed over desks were being labelled “looking "
                "back/up”, and “down” could never fire at all. Confirmed by reading the "
                "library's own source, not guessed.", 10, False, BODY, line=1.20)])
    add_text(s, ML + 0.24, y + 1.76, 5.6, 0.44,
             [PR([R("Solution:  ", 10, True, GREEN, FONT_SB),
                  R("negate at the source + a regression test pinning the convention. "
                    "“down” labels went 0 → 6 on the same images; “back” 2 → 0.",
                    10, False, TEAL_D)], line=1.20)])

    # Problem 6 - the two distraction cases (the user's key ask)
    x2 = ML + 6.05 + 0.28
    w2 = SW - MR - x2
    rect(s, x2, y, w2, 2.30, fill=WHITE, line=BORDER)
    bar(s, x2, y, w2, 0.05, AMBER)
    add_text(s, x2 + 0.24, y + 0.22, w2 - 0.48, 0.24,
             [P("PROBLEM 6 — THE TWO DISTRACTION CASES", 9, True, AMBER, FONT_SB)])
    add_text(s, x2 + 0.24, y + 0.50, w2 - 0.48, 0.30,
             [P("“Head down” means two different things", 13, True, INK, FONT_SB)])
    sub = (w2 - 0.62) / 2
    cases = [("head_down_WITH_device", GREEN,
              "Gaze down + a phone detected overlapping the student. Defensible "
              "behavioural reading → we do flag this."),
             ("head_down_NO_device", AMBER,
              "Gaze down, nothing detected. Could be reading, writing, thinking — or "
              "disengaged. We refuse to guess.")]
    for i, (name, c, desc) in enumerate(cases):
        xx = x2 + 0.24 + i * (sub + 0.14)
        rect(s, xx, y + 0.88, sub, 1.30, fill=PANEL, line=BORDER)
        add_text(s, xx + 0.14, y + 1.00, sub - 0.28, 0.36,
                 [P(name, 8.5, True, c, MONO, line=1.10)])
        add_text(s, xx + 0.14, y + 1.40, sub - 0.28, 0.72,
                 [P(desc, 9, False, BODY, line=1.18)])

    # Bottom: the principle + peer interaction problem
    yy = y + 2.56
    rect(s, ML, yy, 6.05, 1.62, fill=INK, line=None)
    add_text(s, ML + 0.26, yy + 0.20, 5.55, 0.28,
             [P("THE PRINCIPLE WE ADOPTED", 9.5, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])
    add_text(s, ML + 0.26, yy + 0.54, 5.55, 0.94,
             [P("Keep ambiguous cases ambiguous. BOSS — a validated human observation "
                "instrument — has the same unresolved confusion between “quiet, head down, "
                "working” and “quiet, head down, disengaged”. Forcing a confident split "
                "would be less honest than the established instrument itself.",
                10.5, False, WHITE, line=1.24)])

    card(s, x2, yy, w2, 1.62,
         heading="PROBLEM 7 — Peer interaction: a false positive we found and kept",
         lines=["Our geometric pair detector flagged its highest-confidence pair on real "
                "footage. We rendered it and looked: two students at different, "
                "non-adjacent desks, both bent over their own work, not interacting.",
                "We documented the false positive in the code and measured that the "
                "proximity threshold is looser than real desk spacing — rather than "
                "quietly tuning it until the demo looked good."],
         accent=AMBER, heading_size=11.5, body_size=9.5, fill=PANEL)
    return s


def s18_novelty(prs, page):
    s, y = chrome(prs, "Originality",
                  "Deriving what attention means, instead of assuming it",
                  "Everything else in this deck is engineering. This slide is the one "
                  "claim we make that the literature does not.", page)
    items = [
        ("01", "Measure the room before judging the student",
         "Published engagement work assumes students face a teacher, so “facing away” "
         "means distracted. We fit each student's shoulder rays and find where they "
         "converge. If the focus lands among the students it is group work, and facing "
         "away from the front is now correct behaviour, not disengagement."),
        ("02", "The correction is large and has the sign reversed",
         "Replacing camera-relative gaze with room-relative orientation moved the "
         "attentive count from 83 to 1,364 and looking-away from 3,013 to 1,427 on the "
         "same footage. A system without this does not get the number slightly wrong — "
         "it calls an engaged group-work class disengaged."),
        ("03", "Say “we don't know” instead of guessing",
         "The rays need three people to fit; with fewer, the layout reads unknown and "
         "the second score stays blank rather than being invented. 14 of 42 identities "
         "were refused for the same reason, each with the count behind the refusal."),
    ]
    cwid = (CW - 2 * 0.26) / 3
    for i, (num, title, body) in enumerate(items):
        x = ML + i * (cwid + 0.26)
        rect(s, x, y, cwid, 3.20, fill=WHITE, line=BORDER)
        add_text(s, x + 0.24, y + 0.24, cwid - 0.48, 0.5,
                 [P(num, 26, True, RGBColor(0xC9, 0xDA, 0xE2), FONT_SB)])
        add_text(s, x + 0.24, y + 0.84, cwid - 0.48, 0.90,
                 [P(title, 13, True, INK, FONT_SB, line=1.14)])
        add_text(s, x + 0.24, y + 1.78, cwid - 0.48, 1.30,
                 [P(body, 10, False, BODY, line=1.22)])

    yy = y + 3.46
    rect(s, ML, yy, CW, 1.02, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.28, yy + 0.18, CW - 0.56, 0.70,
             [PR([R("Framed honestly:  ", 11, True, TEAL_D, FONT_SB),
                  R("the geometry is elementary — a least-squares intersection of rays. "
                    "What is new is applying it to decide what the engagement score is "
                    "allowed to mean, rather than fixing that in advance. On the 60-clip "
                    "external recording it found 16 group-work and 12 lecture students in "
                    "footage we did not produce.", 11, False, INK)], line=1.24)])
    return s




def s19_scale_bug(prs, page):
    """One wrong constant, found by measuring instead of arguing with labels."""
    s, y = chrome(prs, "Criterion 5 — Problems Faced (3 / 3)",
                  "One constant, tuned for one camera, silently broke every other one",
                  "Found by printing the numbers the labels were derived from, instead "
                  "of arguing with the labels.", page)

    data = [
        ["Inference size", "Persons found", "Box size", "What it means"],
        ["640 px  (1.0x)", "1  (conf 0.92)", "460 x 270", "correct"],
        ["1280 px (2.0x)", "1  (conf 0.61)", "365 x 270", "box shrinking"],
        ["1600 px (2.5x)", "1  (conf 0.46)", "286 x 260", "part of a person"],
        ["1920 px (3.0x)", "0", "-", "person gone"],
    ]
    tw = 7.30
    tbl = table(s, ML, y, tw, [1.85, 1.75, 1.45, 2.25], data,
                row_h=0.40, head_h=0.38, size=9.5, head_size=9.5, col_bold={0})
    for r in range(1, len(data)):
        for c in (0, 1, 2):
            run = tbl.cell(r, c).text_frame.paragraphs[0].runs[0]
            run.font.name = MONO
            run.font.size = Pt(9)
        last = tbl.cell(r, 3).text_frame.paragraphs[0].runs[0]
        last.font.bold = True
        last.font.color.rgb = GREEN if r == 1 else (AMBER if r < 4 else RED)

    yy = y + 0.38 + 4 * 0.40 + 0.26
    card(s, ML, yy, tw, 1.66,
         heading="Why one number broke five different things",
         lines=["imgsz 1920 was tuned on classroom images, where students are small and "
                "enlarging the frame brings them into range. A 640x480 webcam gets "
                "enlarged 3x instead, and a person filling the shot stops looking like "
                "one.",
                "Pose, Face Mesh and head pose all key off the person box. With no box "
                "there is no lean, no eye ratio, no mouth opening, no pitch - those "
                "signals were not wrong, they were never computed. Where the box "
                "survived it had shrunk onto a part of the person, which is what put a "
                "bare hand and a held-up sheet of paper in their own boxes."],
         accent=RED, heading_size=11.5, body_size=9.5, fill=PANEL)

    x2 = ML + tw + 0.30
    w2 = SW - MR - x2
    card(s, x2, y, w2, 2.15,
         heading="The fix, and why it costs nothing",
         lines=["The earlier note argued no rule on frame size could reconcile the two "
                "cases. That was right - the rule is on the size of the upscale, not "
                "the frame. Every classroom gain was already reached at 1.5x or less "
                "(1280x720 x 1.5 is exactly 1920), while every close-up failure sits "
                "above it."],
         accent=TEAL, heading_size=11.5, body_size=9.5, fill=PANEL2)

    yy2 = y + 2.30
    rect(s, x2, yy2, w2, 1.48, fill=WHITE, line=BORDER)
    add_text(s, x2 + 0.22, yy2 + 0.16, w2 - 0.44, 1.25,
             [P("MEASURED BOTH WAYS", 8.5, True, MUTE, FONT_SB),
              PR([R("63 classroom images   ", 10, False, BODY, MONO),
                  R("962 -> 964", 10, True, GREEN, MONO)], space_before=8),
              PR([R("640x480 webcam        ", 10, False, BODY, MONO),
                  R("0 -> 1", 10, True, GREEN, MONO)], space_before=4),
              P("The cap keeps the setting that serves the target domain and refuses "
                "only the upscale that erases a close-up.", 9, False, MUTE,
                line=1.18, space_before=8)])

    yy3 = yy2 + 1.62
    rect(s, x2, yy3, w2, 1.20, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, x2 + 0.22, yy3 + 0.14, w2 - 0.44, 1.02,
             [P("THE LESSON WE WOULD DEFEND", 8.5, True, TEAL_D, FONT_SB),
              P("Every complaint about a wrong label was really a question about a "
                "number - and the numbers were computed and thrown away, leaving only "
                "the label to argue with. tools/probe_signals.py now prints them.",
                9.5, False, INK, line=1.20, space_before=6)])
    return s


def s20_demo(prs, page):
    """What the audience is about to watch, and what to watch for."""
    s, y = chrome(prs, "Live Demonstration",
                  "The room verdict changing, in real time",
                  "Register, capture, graph. The part worth watching is the banner and "
                  "the crosshair, not the boxes.", page)

    steps = [
        ("1", "Register", "Each student is enrolled once and keeps a fixed ID. "
                          "Everything measured afterwards is stored under that ID."),
        ("2", "Capture", "Live webcam. Per person: identity, action, expression, "
                         "orientation - and the room-level layout verdict."),
        ("3", "Graph", "Per-student action transitions and a students-to-objects scene "
                       "graph, with the two scores reported separately."),
    ]
    cwid = (CW - 2 * 0.26) / 3
    for i, (num, title, body) in enumerate(steps):
        x = ML + i * (cwid + 0.26)
        rect(s, x, y, cwid, 1.55, fill=WHITE, line=BORDER)
        add_text(s, x + 0.22, y + 0.18, cwid - 0.44, 0.34,
                 [PR([R(num + "   ", 15, True, TEAL, FONT_SB),
                      R(title, 13, True, INK, FONT_SB)])])
        add_text(s, x + 0.22, y + 0.66, cwid - 0.44, 0.80,
                 [P(body, 9.5, False, BODY, line=1.20)])

    yy = y + 1.78
    rect(s, ML, yy, CW, 1.72, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.28, yy + 0.16, CW - 0.56, 1.40,
             [P("WHAT TO WATCH FOR", 9, True, TEAL_D, FONT_SB),
              P("Three of us sit facing the front - the banner reads Lecture and the "
                "focus crosshair sits out ahead of us. We turn to face each other, the "
                "crosshair slides into the middle of the group, the banner flips to "
                "Group work, and the text changes to say that facing away from the "
                "front is now correct behaviour.",
                11, False, INK, line=1.24, space_before=6),
              P("Nothing is reconfigured between those two states. The system re-derives "
                "what attention means from where the shoulders point.",
                11, True, INK, line=1.24, space_before=6)])

    yy2 = yy + 1.92
    half = (CW - 0.30) / 2
    card(s, ML, yy2, half, 1.32,
         heading="Honest limits, stated up front",
         lines=["The rays need three people. With fewer, the banner says so and the "
                "second score stays blank - we would rather show that than a guess.",
                "A frontal camera cannot tell lying down from sitting upright: both "
                "project to almost the same image geometry."],
         accent=AMBER, heading_size=11, body_size=9.5, fill=PANEL)
    card(s, ML + half + 0.30, yy2, half, 1.32,
         heading="If the live demo fails",
         lines=["outputs/final2 holds a complete 60-clip run on external footage we did "
                "not record: 42 identities, 28 accepted as students, 16 group-work and "
                "12 lecture, 12,789 action classifications.",
                "The report is a static HTML file - no camera, no network, no GPU."],
         accent=GREEN, heading_size=11, body_size=9.5, fill=PANEL)
    return s



# --------------------------------------------------------------------------- #

def main():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs)
    s02_roadmap(prs, 2)
    s03_problem(prs, 3)
    s04_motivation(prs, 4)
    s05_objectives(prs, 5)
    s06_papers(prs, 6)
    s07_lit_technical(prs, 7)
    s08_lit_behavioural(prs, 8)
    s09_gap(prs, 9)
    s10_dataset_primary(prs, 10)
    s12_architecture(prs, 11)
    s14_results(prs, 12)
    s15_problems_perception(prs, 13)
    s16_problems_interpretation(prs, 14)
    s19_scale_bug(prs, 15)
    s18_novelty(prs, 16)
    s20_demo(prs, 17)

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()

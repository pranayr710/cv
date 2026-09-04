"""Build the model-detail deck: what each model is, and what it actually does here.

Every number is measured from this machine by tools/model_facts.py -- parameter
counts read off the loaded graphs, file sizes off disk, and the fine-tuning
metrics re-validated from the trained weights rather than copied out of a log.
A panel asking "how many parameters" should get an answer that came from the
model, not from a datasheet.

Run:  python build_models_ppt.py
Out:  ClassGraph_Models.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "ppt_assets"
OUT = ROOT / "ClassGraph_Models.pptx"

SW, SH = 13.333, 7.5
ML = MR = 0.68
CW = SW - ML - MR

INK = RGBColor(0x10, 0x27, 0x3F)
BODY = RGBColor(0x3C, 0x50, 0x66)
MUTE = RGBColor(0x7B, 0x8C, 0x9E)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
TEAL_D = RGBColor(0x08, 0x59, 0x61)
AMBER = RGBColor(0xB4, 0x7A, 0x14)
GREEN = RGBColor(0x1F, 0x6F, 0x50)
RED = RGBColor(0xA9, 0x33, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF3, 0xF7, 0xFA)
PANEL2 = RGBColor(0xE9, 0xF1, 0xF4)
BORDER = RGBColor(0xDA, 0xE3, 0xEB)
RULE = RGBColor(0xE6, 0xEC, 0xF2)

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"
MONO = "Consolas"


def R(t, size=12, bold=False, color=BODY, font=FONT, italic=False):
    return {"t": t, "size": size, "bold": bold, "color": color,
            "font": font, "italic": italic}


def PR(runs, align=None, space_before=None, space_after=None, line=None):
    return {"runs": runs, "align": align, "space_before": space_before,
            "space_after": space_after, "line": line}


def P(text, size=12, bold=False, color=BODY, font=FONT, italic=False,
      align=None, space_before=None, space_after=None, line=None):
    return PR([R(text, size, bold, color, font, italic)], align,
              space_before, space_after, line)


def add_text(slide, x, y, w, h, blocks, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, blk in enumerate(blocks):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if blk.get("align"):
            para.alignment = blk["align"]
        if blk.get("space_before"):
            para.space_before = Pt(blk["space_before"])
        if blk.get("space_after"):
            para.space_after = Pt(blk["space_after"])
        if blk.get("line"):
            para.line_spacing = blk["line"]
        for r in blk["runs"]:
            run = para.add_run()
            run.text = r["t"]
            run.font.size = Pt(r["size"])
            run.font.bold = r["bold"]
            run.font.italic = r.get("italic", False)
            run.font.name = r["font"]
            run.font.color.rgb = r["color"]
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=0.75):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.adjustments[0] = 0.035
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def bar(slide, x, y, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def chrome(prs, eyebrow, title, lead=None, page=None):
    s = new_slide(prs)
    bar(s, 0, 0, SW, 0.055, TEAL)
    add_text(s, ML, 0.40, CW, 0.26,
             [P(eyebrow.upper(), 9.5, True, TEAL, FONT_SB)])
    add_text(s, ML, 0.68, CW, 0.52, [P(title, 25, True, INK, FONT_SB, line=1.04)])
    y = 1.34
    if lead:
        add_text(s, ML, y, CW, 0.44, [P(lead, 11.5, False, MUTE, line=1.22)])
        y += 0.58
    if page:
        add_text(s, SW - MR - 1.0, SH - 0.46, 1.0, 0.22,
                 [P(str(page), 9, False, MUTE, align=PP_ALIGN.RIGHT)])
    return s, y + 0.06


def table(slide, x, y, w, col_w, data, row_h=0.34, head_h=0.36, size=10,
          head_size=10, col_bold=frozenset()):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w),
                                   Inches(head_h + (rows - 1) * row_h))
    tbl = shape.table
    tbl.first_row = True
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, rows):
        tbl.rows[r].height = Inches(row_h)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                INK if r == 0 else (WHITE if r % 2 else PANEL))
            para = cell.text_frame.paragraphs[0]
            para.line_spacing = 1.0
            for run in para.runs:
                run.font.size = Pt(head_size if r == 0 else size)
                run.font.name = FONT_SB if (r == 0 or c in col_bold) else FONT
                run.font.bold = r == 0 or c in col_bold
                run.font.color.rgb = WHITE if r == 0 else INK
    return tbl


def card(slide, x, y, w, h, heading=None, lines=None, accent=TEAL,
         heading_size=12, body_size=10, fill=PANEL):
    rect(slide, x, y, w, h, fill=fill, line=BORDER)
    bar(slide, x, y, 0.055, h, accent)
    blocks = []
    if heading:
        blocks.append(P(heading, heading_size, True, INK, FONT_SB, line=1.10))
    for ln in (lines or []):
        blocks.append(P(ln, body_size, False, BODY, line=1.22, space_before=6))
    add_text(slide, x + 0.24, y + 0.16, w - 0.44, h - 0.30, blocks)


# --------------------------------------------------------------------------- #
# Measured facts. Regenerate with: python tools/model_facts.py
# --------------------------------------------------------------------------- #

MODELS = [
    # name, params, file MB, purpose
    ("YOLO11m", "20,114,688", "40.7", "Find every person and held object"),
    ("YOLO11m fine-tuned", "20,056,092", "40.5", "Classroom behaviour, 4 classes"),
    ("SCRFD det_10g", "4,225,835", "16.9", "Find faces inside each person box"),
    ("ArcFace w600k_r50", "43,590,976", "174.4", "512-d identity embedding"),
    ("SixDRepNet", "39,323,398", "157.3", "Head pose: yaw, pitch, roll"),
    ("EfficientNet-B0 (emotion)", "3,996,789", "16.0", "Facial expression, 8 classes"),
    ("MediaPipe Pose (BlazePose)", "not exposed", "6.4", "33 body keypoints"),
    ("MediaPipe Face Mesh", "not exposed", "1.2", "468 face landmarks"),
]

PER_CLASS = [
    ["Class", "Precision", "Recall", "mAP@50", "mAP@50-95", "Train boxes"],
    ["write", "0.723", "0.681", "0.767", "0.331", "1,121"],
    ["using_device", "0.749", "0.700", "0.720", "0.336", "2,204"],
    ["sleep", "0.670", "0.495", "0.521", "0.263", "1,512"],
    ["read", "0.455", "0.469", "0.419", "0.143", "1,254"],
    ["all", "0.649", "0.586", "0.607", "0.268", "6,091"],
]


def s01_title(prs):
    s = new_slide(prs)
    bar(s, 0, 0, SW, 0.09, TEAL)
    add_text(s, ML, 2.05, CW, 0.3,
             [P("MODEL CARD DECK", 11, True, TEAL, FONT_SB)])
    add_text(s, ML, 2.44, CW, 1.1,
             [P("The eight models, and what each one is actually for",
                40, True, INK, FONT_SB, line=1.02)])
    add_text(s, ML, 3.80, CW * 0.78, 0.9,
             [P("Purpose, backbone, size, parameter count, training data, output, and "
                "the one place each model is allowed to be wrong. Every number here was "
                "read off the loaded model on this machine, not copied from a datasheet.",
                14, False, BODY, line=1.32)])
    bar(s, ML, 5.05, 1.5, 0.03, TEAL)
    add_text(s, ML, 5.30, CW, 0.9,
             [P("ClassGraph — classroom engagement from video", 13, True, INK, FONT_SB),
              P("Seven pre-trained models used as-is, one fine-tuned on classroom "
                "footage. Regenerate every figure with tools/model_facts.py",
                11, False, MUTE, line=1.26, space_before=6)])
    return s


def s02_overview(prs, page):
    s, y = chrome(prs, "The stack at a glance",
                  "Eight models, 131 million parameters in the pipeline, one of them ours",
                  "131.3 M is the six models whose parameters are countable; MediaPipe "
                  "ships TFLite graphs that do not expose a count. Each model does one "
                  "narrow job it was built for.", page)
    data = [["Model", "Parameters", "MB", "What it is for"]]
    for name, params, mb, purpose in MODELS:
        data.append([name, params, mb, purpose])
    tbl = table(s, ML, y, CW, [3.05, 2.05, 1.05, 5.82], data,
                row_h=0.36, head_h=0.36, size=10.5, head_size=10.5, col_bold={0})
    for r in range(1, len(data)):
        for c in (1, 2):
            run = tbl.cell(r, c).text_frame.paragraphs[0].runs[0]
            run.font.name = MONO
            run.font.size = Pt(10)
        if data[r][0] == "YOLO11m fine-tuned":
            for c in range(4):
                run = tbl.cell(r, c).text_frame.paragraphs[0].runs[0]
                run.font.color.rgb = TEAL_D
                run.font.bold = True

    yy = y + 0.36 + len(MODELS) * 0.36 + 0.24
    half = (CW - 0.30) / 2
    card(s, ML, yy, half, 1.22,
         heading="Why so many small models instead of one big one",
         lines=["A 4 M-parameter face detector and a 20 M-parameter object detector "
                "each solve a problem they were trained for. One network asked to do "
                "both would need far more data than we have, and would fail in ways we "
                "could not attribute to a stage."],
         accent=TEAL, heading_size=11, body_size=9.5, fill=PANEL)
    card(s, ML + half + 0.30, yy, half, 1.22,
         heading="Only one model was trained by us",
         lines=["Seven are used exactly as published. The eighth is YOLO11m fine-tuned "
                "on classroom behaviour, because no public model predicts "
                "read / write / sleep / using_device from a classroom camera. "
                "That is the next four slides."],
         accent=GREEN, heading_size=11, body_size=9.5, fill=PANEL2)
    return s


def s03_detection(prs, page):
    s, y = chrome(prs, "Model 1 of 8 — Detection",
                  "YOLO11m — find every person and every held object",
                  "The first stage. Everything downstream is cropped from the boxes "
                  "this model produces, so its failures are the only ones that cannot "
                  "be recovered later.", page)

    half = (CW - 0.30) / 2
    card(s, ML, y, half, 2.05,
         heading="What it is",
         lines=["Single-stage anchor-free detector, CSPDarknet-style backbone with a "
                "PAN-FPN neck and a decoupled detection head.",
                "20,114,688 parameters, 40.7 MB. Trained by Ultralytics on COCO — "
                "118,000 images, 80 object classes.",
                "Used exactly as published. No fine-tuning: COCO already contains "
                "person, cell phone, laptop, book and bottle, which is what we need."],
         accent=TEAL, heading_size=12, body_size=10, fill=PANEL)

    card(s, ML + half + 0.30, y, half, 2.05,
         heading="What it outputs, and how we use it",
         lines=["Per detection: a box, a class from the 80, and a confidence.",
                "Person boxes become the crop for every other model. Object boxes are "
                "assigned to exactly one student by largest overlap share — before that "
                "rule, 35% of phones were credited to more than one student.",
                "person_conf is 0.30, not the 0.40 default: a back-row student scores "
                "low, and missing one costs more than a spurious box we can filter."],
         accent=TEAL, heading_size=12, body_size=10, fill=PANEL2)

    yy = y + 2.25
    card(s, ML, yy, CW, 1.30,
         heading="The limitation a panel should know about — COCO has 80 classes, not 500",
         lines=["A spectacle case, a poster and a water bottle are not the same problem: "
                "the bottle is a COCO class and is detected, the other two are not and "
                "never will be by this model. We report what the model can name and stay "
                "silent otherwise, rather than guessing a label from shape.",
                "The classes we actually consume are person, cell phone, laptop, book, "
                "bottle, cup and keyboard. A held object outside that set is evidence "
                "that the hands are busy, which the action rules use without naming it."],
         accent=AMBER, heading_size=11.5, body_size=9.5, fill=PANEL)

    yy2 = yy + 1.50
    rect(s, ML, yy2, CW, 1.28, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.28, yy2 + 0.16, CW - 0.56, 1.00,
             [P("THE SETTING THAT MATTERED MOST", 9, True, TEAL_D, FONT_SB),
              P("Inference size is not a free parameter. imgsz 1920 suits a classroom "
                "camera, where students are small and enlarging the frame helps. On a "
                "640x480 webcam it enlarges 3x and the detector finds nobody at all — "
                "so the upscale is now capped at 1.5x, which costs nothing on classroom "
                "footage (962 -> 964 persons) and takes the webcam from 0 to 1.",
                11, False, INK, line=1.24, space_before=6)])
    return s


def s04_finetune_why(prs, page):
    s, y = chrome(prs, "Model 2 of 8 — Fine-tuning (1 / 3)",
                  "Why we had to train one model ourselves",
                  "Everything else in this deck is a model somebody else trained. This "
                  "is the one gap no public checkpoint filled.", page)

    third = (CW - 2 * 0.26) / 3
    blocks = [
        ("THE GAP", "COCO knows objects, not behaviour",
         "COCO can tell us a laptop is present. It cannot tell us whether the student "
         "is typing on it, reading from it, or asleep beside it. No public detector "
         "predicts classroom behaviour classes from a room camera."),
        ("THE CHOICE", "Fine-tune, do not train from scratch",
         "935 labelled images cannot train a 20 M-parameter detector from random "
         "weights. Starting from COCO weights means the backbone already knows edges, "
         "texture and human shape, and only the class semantics have to be learned."),
        ("THE SCOPE", "Four classes, deliberately few",
         "read, write, sleep, using_device. Each is a visible posture-and-object "
         "configuration. We did not add classes like 'bored' or 'confused' — those are "
         "interpretations, not things a camera can see."),
    ]
    for i, (tag, title, body) in enumerate(blocks):
        x = ML + i * (third + 0.26)
        rect(s, x, y, third, 2.55, fill=WHITE, line=BORDER)
        bar(s, x, y, third, 0.05, TEAL)
        add_text(s, x + 0.24, y + 0.22, third - 0.48, 0.24,
                 [P(tag, 9, True, TEAL, FONT_SB)])
        add_text(s, x + 0.24, y + 0.54, third - 0.48, 0.56,
                 [P(title, 13, True, INK, FONT_SB, line=1.10)])
        add_text(s, x + 0.24, y + 1.18, third - 0.48, 1.20,
                 [P(body, 10, False, BODY, line=1.22)])

    yy = y + 2.78
    card(s, ML, yy, CW, 1.50,
         heading="Transfer learning, stated precisely",
         lines=["A detector is a backbone that turns pixels into features, a neck that "
                "mixes them across scales, and a head that turns features into boxes and "
                "class scores. Only the head's meaning is task-specific; edges and human "
                "shape are the same in COCO and in a classroom.",
                "So we keep the COCO weights as the starting point and let gradient "
                "descent move all of them. The backbone barely moves because it is "
                "already right; the head moves a lot, because 80 COCO classes have to "
                "become 4 behaviour classes."],
         accent=GREEN, heading_size=12, body_size=10, fill=PANEL2)
    return s


def s05_finetune_how(prs, page):
    s, y = chrome(prs, "Model 2 of 8 — Fine-tuning (2 / 3)",
                  "Exactly what was trained, on what, with which settings",
                  "Read straight out of runs/behaviour/merged4_aug/args.yaml — the "
                  "configuration the run actually used, not one written down afterwards.",
                  page)

    half = (CW - 0.30) / 2
    cfg = [
        ["Setting", "Value", "Why"],
        ["Starting weights", "yolo11m.pt", "COCO-pretrained, not random"],
        ["Layers frozen", "none (freeze: null)", "all 20.1 M updated"],
        ["Dataset", "behaviour_merged", "877 train / 58 val images"],
        ["Labelled boxes", "6,091 train", "4 classes"],
        ["Epochs", "42 of 60", "early-stopped, patience 15"],
        ["Best epoch", "27", "later epochs overfit"],
        ["Image size", "640", "1600 needed 9.05 GB on a 6.4 GB card"],
        ["Batch", "8", "largest that fits in VRAM"],
        ["Optimiser / LR", "auto, lr0 0.01", "Ultralytics default"],
        ["Wall-clock", "23 min", "RTX 4050 laptop"],
    ]
    tbl = table(s, ML, y, half, [1.72, 1.85, 2.35], cfg,
                row_h=0.315, head_h=0.33, size=9.5, head_size=9.5, col_bold={0})
    for r in range(1, len(cfg)):
        run = tbl.cell(r, 1).text_frame.paragraphs[0].runs[0]
        run.font.name = MONO
        run.font.size = Pt(9)
        run.font.color.rgb = TEAL_D
        run.font.bold = True

    x2 = ML + half + 0.30
    card(s, x2, y, half, 2.10,
         heading="Which layers were fine-tuned? All of them.",
         lines=["args.yaml records freeze: null, so no layer was held fixed — every one "
                "of the 20.1 M parameters received gradients.",
                "This is the right choice at this scale. Freezing the backbone is for "
                "when the new data is tiny or very close to the original domain; a "
                "classroom seen from a high rear corner is neither. The COCO weights are "
                "the starting point, not a fixed feature extractor.",
                "The parameter count drops slightly, 20,114,688 to 20,056,092, because "
                "the detection head now predicts 4 classes instead of 80."],
         accent=GREEN, heading_size=12, body_size=9.5, fill=PANEL2)

    yy = y + 2.30
    card(s, x2, yy, half, 1.55,
         heading="The dataset, and its honest weakness",
         lines=["935 images total, merged from classroom sets and hand-checked. Train "
                "boxes per class: using_device 2,204, sleep 1,512, read 1,254, "
                "write 1,121.",
                "The validation split is 58 images. That is small enough that a single "
                "hard image moves the metric, and we say so rather than quoting the "
                "number as if it were a benchmark result."],
         accent=AMBER, heading_size=12, body_size=9.5, fill=PANEL)
    return s


def s06_finetune_result(prs, page):
    s, y = chrome(prs, "Model 2 of 8 — Fine-tuning (3 / 3)",
                  "What fine-tuning bought, and where it still fails",
                  "Re-validated from the trained weights for this deck, not copied from "
                  "the training log.", page)

    tw = 7.55
    tbl = table(s, ML, y, tw, [1.75, 1.20, 1.05, 1.15, 1.35, 1.05], PER_CLASS,
                row_h=0.40, head_h=0.38, size=10, head_size=9.5, col_bold={0})
    for r in range(1, len(PER_CLASS)):
        for c in range(1, 6):
            run = tbl.cell(r, c).text_frame.paragraphs[0].runs[0]
            run.font.name = MONO
            run.font.size = Pt(9.5)
        name = PER_CLASS[r][0]
        colour = (INK if name == "all" else
                  GREEN if r <= 2 else (AMBER if r == 3 else RED))
        run = tbl.cell(r, 0).text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = colour
        run.font.bold = True
        tbl.cell(r, 3).text_frame.paragraphs[0].runs[0].font.color.rgb = colour
        tbl.cell(r, 3).text_frame.paragraphs[0].runs[0].font.bold = True

    yy = y + 0.38 + len(PER_CLASS) * 0.40 + 0.24
    card(s, ML, yy, tw, 1.42,
         heading="Did fine-tuning improve accuracy? Yes — measurably.",
         lines=["mAP@50 on the validation split rose from 0.188 after the first epoch "
                "to 0.607 at the best epoch: a 3.2x improvement as the head learned the "
                "four classes. The COCO model scores zero on this task by construction, "
                "because it cannot emit these labels at all.",
                "Early stopping at epoch 42 with the best at 27 is the useful detail: "
                "the model stopped improving on unseen images long before it stopped "
                "improving on the training set."],
         accent=GREEN, heading_size=11.5, body_size=9.5, fill=PANEL2)

    x2 = ML + tw + 0.30
    w2 = SW - MR - x2
    card(s, x2, y, w2, 2.35,
         heading="'read' is the weakest class — and we predicted it",
         lines=["read scores mAP@50 0.419 against write at 0.767. That is not noise: "
                "reading and writing differ only by what the hands are doing, and the "
                "closest published work reaches 57.8% on writing even with a strong "
                "temporal model.",
                "So the pipeline does not present them as equals. When a book is visible "
                "but the hands are not, the reported action is 'reading or writing' — one "
                "label covering both — and writing is flagged inferred, never direct."],
         accent=RED, heading_size=12, body_size=9.5, fill=PANEL)

    card(s, x2, y + 2.55, w2, 2.05,
         heading="Where this model sits in the pipeline",
         lines=["It is a second opinion, not the primary signal. Actions are decided "
                "first by geometry — object overlap, wrist position, head pitch — which "
                "is auditable and needs no training data.",
                "The behaviour model supplies a label where geometry is silent. Keeping "
                "it subordinate is why a 0.607 mAP model is safe to include: its "
                "mistakes are visible as a disagreement, not as the only answer."],
         accent=TEAL, heading_size=12, body_size=9.5, fill=PANEL2)
    return s


def s07_identity(prs, page):
    s, y = chrome(prs, "Models 3 and 4 of 8 — Identity",
                  "SCRFD finds the face, ArcFace says whose it is",
                  "Identity is what makes a per-frame observation into a student's "
                  "history. Both models are used exactly as published.", page)

    half = (CW - 0.30) / 2
    card(s, ML, y, half, 2.55,
         heading="SCRFD det_10g — face detection",
         lines=["4,225,835 parameters, 16.9 MB, ONNX on GPU. Part of InsightFace's "
                "buffalo_l pack, trained on WIDER FACE.",
                "Sample-redistribution single-stage detector: it spends most of its "
                "compute on the small-face scales, which is exactly the classroom case.",
                "Output per face: a box, a confidence, and 5 keypoints (both eyes, nose, "
                "two mouth corners). The keypoints matter — they align the crop before "
                "recognition and before expression."],
         accent=TEAL, heading_size=12, body_size=9.5, fill=PANEL)

    card(s, ML + half + 0.30, y, half, 2.55,
         heading="ArcFace w600k_r50 — identity embedding",
         lines=["43,590,976 parameters, 174.4 MB — the largest model in the stack. "
                "ResNet-50 backbone, trained on WebFace600K.",
                "Output: one 512-dimensional unit vector per face. Not a name and not a "
                "class — a point on a hypersphere, where the same person lands close "
                "together and different people land far apart.",
                "The additive angular margin in its training loss is what forces that "
                "separation, which is why cosine distance between two embeddings is a "
                "meaningful identity score at all."],
         accent=TEAL, heading_size=12, body_size=9.5, fill=PANEL2)

    yy = y + 2.75
    card(s, ML, yy, CW, 1.85,
         heading="What we build on top — the part that is ours",
         lines=["A raw embedding match is not enough in a classroom. Two students can "
                "look alike, and the same student's embedding drifts as they turn. So "
                "identity is resolved by constrained agglomerative clustering over "
                "cosine distance, with a hard cannot-link between tracks that appear in "
                "the same frame — one body cannot be two people at once, and that "
                "constraint is free information the embedding does not carry.",
                "On registration each student is enrolled once and keeps a fixed ID. "
                "98.6% of detections receive an ID; where no face is good enough, the "
                "detection is reported as unidentified rather than attributed to the "
                "nearest match — 6 of 42 identities in the 60-clip run were refused on "
                "exactly this ground."],
         accent=GREEN, heading_size=12, body_size=9.5, fill=PANEL)
    return s


def s08_geometry(prs, page):
    s, y = chrome(prs, "Models 5, 6 and 7 of 8 — Geometry",
                  "Head pose and landmarks: the measurements behind every action",
                  "These three produce no labels at all. They produce numbers, and the "
                  "action rules read those numbers.", page)

    third = (CW - 2 * 0.26) / 3
    items = [
        ("SixDRepNet", TEAL,
         ["39,323,398 parameters, 157 MB. RepVGG backbone. Checkpoint 6DRepNet_300W_LP_AFLW2000 — trained on 300W-LP, evaluated on AFLW2000.",
          "Predicts head rotation as a 6-D continuous representation rather than three "
          "Euler angles — Euler angles are discontinuous at their wrap-around, which "
          "makes them a poor regression target.",
          "Output: yaw, pitch, roll in degrees. We use pitch to separate a bowed head "
          "from a raised one. Yaw is now a fallback only, because it is measured "
          "relative to the camera and every seat has a different angle to it."]),
        ("MediaPipe Pose", TEAL,
         ["BlazePose, 6.4 MB TFLite, CPU. model_complexity 1, so the \"full\" landmark model. Two-stage: a detector, then a landmark model.",
          "Output: 33 body keypoints with visibility scores — shoulders, hips, wrists, "
          "nose.",
          "Two things depend on it. Wrist position relative to the face and the desk "
          "gives raised hand, head-on-hand and hands-low. Shoulder direction gives "
          "facing_direction in image space, which is the input to the room-layout "
          "measurement and carries no camera constant."]),
        ("MediaPipe Face Mesh", TEAL,
         ["1.2 MB TFLite, CPU. Its own detector is bypassed — it runs on the crop SCRFD "
          "already found, so the two never disagree about where the face is.",
          "Output: 478 landmarks with refinement on; we keep the canonical 468 and drop the 10 iris points, because the frozen output schema says 468. Refinement still sharpens the eye points the EAR depends on.",
          "We derive two scalars: eye aspect ratio for closed eyes, and mouth opening "
          "ratio for a yawn. Both are ratios of distances between landmarks, so they are "
          "scale-invariant — a face near the camera and one far away give comparable "
          "numbers."]),
    ]
    for i, (title, accent, lines) in enumerate(items):
        x = ML + i * (third + 0.26)
        rect(s, x, y, third, 3.55, fill=WHITE, line=BORDER)
        bar(s, x, y, third, 0.05, accent)
        add_text(s, x + 0.24, y + 0.24, third - 0.48, 0.34,
                 [P(title, 14, True, INK, FONT_SB)])
        blocks = [P(ln, 9.5, False, BODY, line=1.22, space_before=7) for ln in lines]
        add_text(s, x + 0.24, y + 0.68, third - 0.48, 2.70, blocks)

    yy = y + 3.75
    rect(s, ML, yy, CW, 1.10, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.28, yy + 0.15, CW - 0.56, 0.85,
             [P("WHY THESE ARE THE MOST IMPORTANT MODELS IN THE STACK", 9, True,
                TEAL_D, FONT_SB),
              P("A classifier gives a label you must trust. These give a measurement you "
                "can check. When a threshold is wrong, the number shows it — which is how "
                "the yawn threshold was found to be unreachable (set at 0.55, the signal "
                "never exceeded 0.274) and how the eye ratio was found to be saturated.",
                11, False, INK, line=1.24, space_before=6)])
    return s


def s09_expression(prs, page):
    s, y = chrome(prs, "Model 8 of 8 — Expression",
                  "EfficientNet-B0, and why we report 3 labels from its 8",
                  "The smallest model in the stack, and the one where we most "
                  "deliberately discard information.", page)

    half = (CW - 0.30) / 2
    card(s, ML, y, half, 2.45,
         heading="What it is",
         lines=["EfficientNet-B0 backbone, 3,996,789 parameters, 16.0 MB ONNX. Run on "
                "CPU on purpose: YOLO and SixDRepNet already occupy the 6.4 GB card.",
                "Checkpoint enet_b0_8_best_vgaf from EmotiEffLib — trained on AffectNet "
                "and validated on VGAF, an in-the-wild video group-affect set, which is "
                "much closer to a classroom than a posed studio dataset.",
                "Input 224x224. Output: 8 AffectNet class probabilities."],
         accent=TEAL, heading_size=12, body_size=9.5, fill=PANEL)

    card(s, ML + half + 0.30, y, half, 2.45,
         heading="The 8 to 3 mapping, and why it is not lossy in the file",
         lines=["Reported: happy, sad, neutral. Anger, Contempt, Disgust, Fear and "
                "Surprise all map to neutral.",
                "They map to neutral rather than to sad on purpose. Folding anger into "
                "sadness would assert something the model never said; mapping it to "
                "neutral says only 'not one of the three we report', which is true.",
                "The full 8-class distribution is kept in the output and the mapping "
                "lives in config, so the collapse is auditable and reversible — a "
                "reviewer can ask for the other five and get them."],
         accent=GREEN, heading_size=12, body_size=9.5, fill=PANEL2)

    yy = y + 2.65
    card(s, ML, yy, CW, 1.65,
         heading="Alignment: a measured decision, not a default",
         lines=["AffectNet was trained on aligned faces, so feeding a raw box crop puts "
                "the model out of distribution. We align each crop using SCRFD's five "
                "keypoints before classifying.",
                "On the 60-clip external run this model produced 9,321 classifications — "
                "neutral 4,011, happy 707, sad 665 — and 3,468 came back 'uncertain'. "
                "Those are faces too small or too turned to classify, and they are "
                "reported as uncertain rather than defaulted to neutral, which would "
                "have silently inflated the neutral count by 37%."],
         accent=AMBER, heading_size=12, body_size=9.5, fill=PANEL)
    return s


def s10_summary(prs, page):
    s, y = chrome(prs, "Summary",
                  "What each model is allowed to be wrong about",
                  "The question a panel is really asking is not how big each model is, "
                  "but what happens when it fails.", page)

    data = [
        ["Model", "Output", "Failure mode", "How the system contains it"],
        ["YOLO11m", "boxes + 80 classes", "misses a small student",
         "conf floor 0.30; upscale capped at 1.5x"],
        ["YOLO11m fine-tuned", "4 behaviour classes", "confuses read with write",
         "reported as 'reading or writing', marked inferred"],
        ["SCRFD", "face box + 5 points", "no face when head is down",
         "body pose recovers the student; else unknown"],
        ["ArcFace", "512-d embedding", "two students look alike",
         "cannot-link between co-occurring tracks"],
        ["SixDRepNet", "yaw, pitch, roll", "yaw is camera-relative",
         "demoted to fallback; room layout used instead"],
        ["MediaPipe Pose", "33 keypoints", "fits a skeleton to furniture",
         "requires both shoulders before a lean is reported"],
        ["MediaPipe Face Mesh", "468 landmarks", "eye ratio saturates",
         "known open — needs open-eye calibration"],
        ["EfficientNet-B0", "8 emotion classes", "small or turned faces",
         "reported 'uncertain', never defaulted to neutral"],
    ]
    tbl = table(s, ML, y, CW, [2.35, 2.05, 2.55, 5.02], data,
                row_h=0.40, head_h=0.38, size=9.5, head_size=9.5, col_bold={0})
    for r in range(1, len(data)):
        tbl.cell(r, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RED
        run = tbl.cell(r, 3).text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = GREEN
        run.font.bold = True

    yy = y + 0.38 + 8 * 0.40 + 0.26
    rect(s, ML, yy, CW, 1.10, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.28, yy + 0.15, CW - 0.56, 0.85,
             [P("THE ONE ROW WE WOULD DEFEND HARDEST", 9, True, TEAL_D, FONT_SB),
              P("Every row's containment strategy ends in the same place: when a model "
                "cannot answer, the system says so instead of substituting its best "
                "guess. 14 of 42 identities in the external run were refused on those "
                "grounds, each with the count behind the refusal recorded.",
                11, False, INK, line=1.24, space_before=6)])
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs)
    s02_overview(prs, 2)
    s03_detection(prs, 3)
    s04_finetune_why(prs, 4)
    s05_finetune_how(prs, 5)
    s06_finetune_result(prs, 6)
    s07_identity(prs, 7)
    s08_geometry(prs, 8)
    s09_expression(prs, 9)
    s10_summary(prs, 10)

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
